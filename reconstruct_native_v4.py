from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# COLORS
BEIGE_BG = RGBColor(240, 242, 235)
DARK_GREEN = RGBColor(27, 43, 17)
LIGHT_GREEN_PANEL = RGBColor(225, 234, 216)
GOLD = RGBColor(181, 150, 75)
WHITE = RGBColor(255, 255, 255)

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_native_stat_card(slide, left, top, width, height, number, label):
    # Shadow (Simplified by adding a slightly offset dark shape behind)
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.05), top + Inches(0.05), width, height)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(200, 200, 200)
    shadow.line.fill.background()
    
    # Main Box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = DARK_GREEN
    shape.line.width = Pt(4)
    
    # Number (Gold)
    num_box = slide.shapes.add_textbox(left, top + Inches(0.3), width, Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.name = 'Playfair Display'
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Label
    lbl_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(1.3), width - Inches(0.2), height - Inches(1.4))
    tf = lbl_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = 'Outfit'
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER

def add_slide_2_native(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE_BG)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1))
    title.text_frame.text = "63 Million Invisible Workers. $1 Trillion in Unpaid Labor."
    p = title.text_frame.paragraphs[0]
    p.font.name = 'Playfair Display'
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(15), Inches(0.4))
    sub.text_frame.text = "Family caregivers are the backbone of American healthcare — yet they remain completely invisible."
    p = sub.text_frame.paragraphs[0]
    p.font.name = 'Outfit'
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER

    # Top Row Cards
    card_w = Inches(3.5)
    card_h = Inches(3)
    spacing = Inches(0.3)
    start_x = (prs.slide_width - (4 * card_w + 3 * spacing)) / 2
    
    add_native_stat_card(slide, start_x, Inches(1.8), card_w, card_h, "63M", "Unpaid family caregivers in the US, 19% increase since 2020.")
    add_native_stat_card(slide, start_x + (card_w + spacing), Inches(1.8), card_w, card_h, "$1 Trillion", "Annual economic value of unpaid caregiving in 2024.")
    add_native_stat_card(slide, start_x + 2*(card_w + spacing), Inches(1.8), card_w, card_h, "27 hrs/wk", "Average caregiving hours; 40%+ provide high-intensity care.")
    add_native_stat_card(slide, start_x + 3*(card_w + spacing), Inches(1.8), card_w, card_h, "$10,200/y", "Average out-of-pocket spending per caregiver annually.")

    # Bottom Panel
    panel_h = Inches(3.5)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.2), Inches(15), panel_h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT_GREEN_PANEL
    panel.line.color.rgb = DARK_GREEN
    panel.line.width = Pt(1.5)
    
    # Panel Title
    p_title = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(15), Inches(0.5))
    p_title.text_frame.text = "Why has no one captured this market?"
    p = p_title.text_frame.paragraphs[0]
    p.font.name = 'Playfair Display'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    # Quadrants (Simplified as text boxes)
    quads = [
        ("Invisible to Tech", "Caregivers are the operating system keeping care together, yet tech ignores their workflow."),
        ("The Hands-Full Problem", "Active care requires full physical involvement. Typing is physically impossible."),
        ("The Medical Knowledge Gap", "78% manage complex tasks with zero clinical background."),
        ("Fragmented Tools", "Single-dimension apps solve one problem without clinical context.")
    ]
    
    for i, (head, body) in enumerate(quads):
        col = i % 2
        row = i // 2
        left = Inches(1 + col * 7.5)
        top = Inches(6 + row * 1.2)
        
        box = slide.shapes.add_textbox(left, top, Inches(6.5), Inches(1.2))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = head
        p.font.name = 'Playfair Display'
        p.font.size = Pt(22)
        p.font.bold = True
        
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.name = 'Outfit'
        p2.font.size = Pt(14)
        
    # Footer Banner
    footer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(8.5), Inches(9), Inches(0.6))
    footer.fill.solid()
    footer.fill.fore_color.rgb = DARK_GREEN
    
    ft_text = slide.shapes.add_textbox(Inches(3.5), Inches(8.55), Inches(9), Inches(0.5))
    ft_text.text_frame.text = "By 2030, all 73 million boomers will be 65+."
    p = ft_text.text_frame.paragraphs[0]
    p.font.name = 'Outfit'
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def generate_native_deck():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Slide 1: Title (Skip for brevity in this test, but I'll add it in the final)
    # Slide 2: Benchmark
    add_slide_2_native(prs)
    
    output_path = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/FINAL_Pitchdeck/AuraGuide_Native_Reconstruction_TEST.pptx"
    prs.save(output_path)
    print(f"Native Test Deck saved to {output_path}")

if __name__ == "__main__":
    generate_native_deck()
