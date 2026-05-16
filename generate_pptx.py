from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Colors
BEIGE = RGBColor(240, 242, 235)
DARK_GREEN = RGBColor(27, 43, 17)
LIME = RGBColor(204, 255, 0)

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title_text, subtitle_text, image_path=None):
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BEIGE)
    
    if image_path and os.path.exists(image_path):
        # Center the cover image or fill the slide
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=prs.slide_width)
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(1.5))
    tf = title.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.name = 'Playfair Display'
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(1))
    tf = subtitle.text_frame
    tf.text = subtitle_text
    p = tf.paragraphs[0]
    p.font.name = 'Outfit'
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BEIGE)
    
    # Side Border
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_GREEN
    shape.line.fill.background()
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    tf = title.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.name = 'Playfair Display'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    
    # Content
    content = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(14), Inches(6))
    tf = content.text_frame
    tf.word_wrap = True
    
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.name = 'Outfit'
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_GREEN
        p.space_after = Pt(18)

def generate_deck():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    cover_img = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/pitch-deck/assets/cover.png"
    
    # Slide 1: Cover
    add_title_slide(prs, "AURAPATH", "Intelligent Eyewear for Cognitive Independence.", cover_img)
    
    # Slide 2: The Problem
    add_content_slide(prs, "Memory Loss Is a Dignity Crisis.", [
        "• 55 million people worldwide live with dementia.",
        "• 60% of dementia patients will wander, leading to trauma.",
        "• Social Isolation: Failed recognition accelerates cognitive decline.",
        "• Medication Danger: Prescription errors kill 7,000+ patients annually.",
        "• Market Size: $360B annual cost of dementia care in the US."
    ])
    
    # Slide 3: The Solution
    add_content_slide(prs, "AuraPath: A Digital Prosthetic for Memory.", [
        "• Social Recognition: Names whispered via bone-conduction audio.",
        "• Medication Safety: Computer vision verifies labels before every dose.",
        "• Wandering Prevention: Geofencing and haptic navigation guide users home.",
        "• Privacy-First: All inference runs on-device via Apple MLX NPUs.",
        "• No Stigma: Indistinguishable from normal eyewear."
    ])
    
    # Slide 4: Market Timing
    add_content_slide(prs, "Three Forces. One Inflection Point.", [
        "• Edge AI Power: Apple M-series NPUs enable local vision models (200ms latency).",
        "• Aging Population: 10,000 boomers turn 65 daily; cases tripling by 2050.",
        "• Category Gap: Passive aids react after failure; AuraPath prevents it proactively."
    ])

    # Slide 5: The Ask
    add_content_slide(prs, "Raising $750K Pre-Seed.", [
        "• 40% Hardware R&D: MLX optimization and NPU feasibility.",
        "• 25% Regulatory & IP: FDA classification and ISO 13485 QMS.",
        "• 20% Clinical Pilot: Caregiver beta and facility partnerships.",
        "• 15% Operations: Advisory board and team scaling.",
        "• Terms: Post-Money SAFE | $4M–$6M Cap | 20% Discount"
    ])
    
    output_path = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/AuraPath_PitchDeck_v2.pptx"
    prs.save(output_path)
    print(f"Deck saved to {output_path}")

if __name__ == "__main__":
    generate_deck()
