from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
BEIGE = RGBColor(240, 242, 235)
DARK_GREEN = RGBColor(27, 43, 17)
LIME = RGBColor(204, 255, 0)
WHITE = RGBColor(255, 255, 255)

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, size=44):
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(15), Inches(1))
    tf = title.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.name = 'Playfair Display'
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    return title

def add_subtitle(slide, text, top=1.2):
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(15), Inches(0.5))
    tf = subtitle.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.name = 'Outfit'
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GREEN
    p.font.italic = True

def add_stat_card(slide, left, top, width, height, number, label):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = DARK_GREEN
    
    num_box = slide.shapes.add_textbox(left, top + Inches(0.2), width, Inches(0.6))
    num_box.text_frame.text = number
    num_box.text_frame.paragraphs[0].font.name = 'Playfair Display'
    num_box.text_frame.paragraphs[0].font.size = Pt(40)
    num_box.text_frame.paragraphs[0].font.bold = True
    num_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    lbl_box = slide.shapes.add_textbox(left, top + Inches(0.8), width, Inches(0.4))
    lbl_box.text_frame.text = label
    lbl_box.text_frame.paragraphs[0].font.name = 'Outfit'
    lbl_box.text_frame.paragraphs[0].font.size = Pt(12)
    lbl_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    add_title(slide, "The Caregiver Collapse Cycle.")
    add_subtitle(slide, "A crisis hiding in plain sight.")
    
    img = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/temp_pdf_analysis/p3.png"
    if os.path.exists(img):
        slide.shapes.add_picture(img, Inches(1), Inches(2), width=Inches(14))

def add_solution_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    add_title(slide, "AuraGuide: Medical Intelligence Partner.")
    
    img = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/temp_pdf_analysis/p4.png"
    if os.path.exists(img):
        slide.shapes.add_picture(img, Inches(1), Inches(2), width=Inches(14))

def add_how_it_works_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    add_title(slide, "How It Works: The OS for Home Care.")
    img = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/temp_pdf_analysis/p5.png"
    if os.path.exists(img): slide.shapes.add_picture(img, Inches(1), Inches(2), width=Inches(14))

def add_conversations_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    add_title(slide, "Conversations That Save Lives.")
    img = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/temp_pdf_analysis/p6.png"
    if os.path.exists(img): slide.shapes.add_picture(img, Inches(0.5), Inches(1.5), width=Inches(15))

def add_market_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    add_title(slide, "A $72B Market Opportunity.")
    img = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/temp_pdf_analysis/p7.png"
    if os.path.exists(img): slide.shapes.add_picture(img, Inches(0.5), Inches(1.5), width=Inches(15))

def add_business_model_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    add_title(slide, "Business Model: Hardware + SaaS.")
    img = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/temp_pdf_analysis/p8.png"
    if os.path.exists(img): slide.shapes.add_picture(img, Inches(1), Inches(2), width=Inches(14))

def add_competition_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    add_title(slide, "Competitive Landscape.")
    img = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/temp_pdf_analysis/p9.png"
    if os.path.exists(img): slide.shapes.add_picture(img, Inches(1), Inches(2), width=Inches(14))

def add_traction_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    add_title(slide, "Traction and Roadmap.")
    img = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/temp_pdf_analysis/p10.png"
    if os.path.exists(img): slide.shapes.add_picture(img, Inches(1), Inches(2), width=Inches(14))

def add_team_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    add_title(slide, "The Team.")
    img = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/temp_pdf_analysis/p11.png"
    if os.path.exists(img): slide.shapes.add_picture(img, Inches(1), Inches(2), width=Inches(14))

def add_ask_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    add_title(slide, "The Ask: $500,000 Pre-Seed.")
    img = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/temp_pdf_analysis/p12.png"
    if os.path.exists(img): slide.shapes.add_picture(img, Inches(1), Inches(2), width=Inches(14))

def add_ecosystem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    add_title(slide, "The Viktron Health Ecosystem.")
    img = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/temp_pdf_analysis/p13.png"
    if os.path.exists(img): slide.shapes.add_picture(img, Inches(1), Inches(2), width=Inches(14))

def generate_full_deck():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    img = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/temp_pdf_analysis/p1.png"
    if os.path.exists(img): slide.shapes.add_picture(img, Inches(0), Inches(0), width=Inches(16))
    
    # 2. Market Scope
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    img = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/temp_pdf_analysis/p2.png"
    if os.path.exists(img): slide.shapes.add_picture(img, Inches(0), Inches(0), width=Inches(16))
    
    add_problem_slide(prs) # 3
    add_solution_slide(prs) # 4
    add_how_it_works_slide(prs) # 5
    add_conversations_slide(prs) # 6
    add_market_slide(prs) # 7
    add_business_model_slide(prs) # 8
    add_competition_slide(prs) # 9
    add_traction_slide(prs) # 10
    add_team_slide(prs) # 11
    add_ask_slide(prs) # 12
    add_ecosystem_slide(prs) # 13
    
    output_path = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/FINAL_Pitchdeck/AuraGuide_PitchDeck_v4_Full_HighFidelity.pptx"
    prs.save(output_path)
    print(f"Full 13-slide high-fidelity deck saved to {output_path}")

if __name__ == "__main__":
    generate_full_deck()
