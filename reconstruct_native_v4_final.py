from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import os

# COLORS
BEIGE_BG = RGBColor(240, 242, 235)
DARK_GREEN = RGBColor(27, 43, 17)
LIGHT_GREEN_PANEL = RGBColor(225, 234, 216)
GOLD = RGBColor(181, 150, 75)
WHITE = RGBColor(255, 255, 255)
ACCENT_LIME = RGBColor(204, 255, 0)

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, size=44, color=DARK_GREEN):
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(15), Inches(1))
    tf = title.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.name = 'Playfair Display'
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return title

def add_native_stat_card(slide, left, top, width, height, number, label):
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = DARK_GREEN
    box.line.width = Pt(4)
    
    n = slide.shapes.add_textbox(left, top + Inches(0.3), width, Inches(1))
    n.text_frame.text = number
    p = n.text_frame.paragraphs[0]
    p.font.name = 'Playfair Display'
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    l = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(1.3), width - Inches(0.2), height - Inches(1.4))
    l.text_frame.text = label
    p = l.text_frame.paragraphs[0]
    p.font.name = 'Outfit'
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER

def add_slide_1_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE_BG)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-4), Inches(-1), Inches(12), Inches(11))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_GREEN
    shape.line.fill.background()
    title = slide.shapes.add_textbox(Inches(8), Inches(3.2), Inches(7), Inches(1.5))
    title.text_frame.text = "AURAGUIDE"
    p = title.text_frame.paragraphs[0]
    p.font.name = 'Playfair Display'
    p.font.size = Pt(100)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    tag = slide.shapes.add_textbox(Inches(8), Inches(4.8), Inches(7), Inches(0.8))
    tag.text_frame.text = "A Lifeline for Family Caregivers."
    p = tag.text_frame.paragraphs[0]
    p.font.name = 'Outfit'
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = DARK_GREEN

def add_slide_2_stats(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE_BG)
    add_title(slide, "63 Million Invisible Workers.")
    add_native_stat_card(slide, Inches(1), Inches(2), Inches(3.2), Inches(3), "63M", "Unpaid Family Caregivers")
    add_native_stat_card(slide, Inches(4.5), Inches(2), Inches(3.2), Inches(3), "$1T", "Economic Value")
    add_native_stat_card(slide, Inches(8), Inches(2), Inches(3.2), Inches(3), "27h/wk", "Avg. Care Hours")
    add_native_stat_card(slide, Inches(11.5), Inches(2), Inches(3.2), Inches(3), "87%", "Report Stress")

def add_slide_7_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE_BG)
    add_title(slide, "A $72B Market at Inflection Point.")
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(2), Inches(6), Inches(6))
    c1.fill.solid()
    c1.fill.fore_color.rgb = LIGHT_GREEN_PANEL
    c1.line.color.rgb = DARK_GREEN
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.8), Inches(2.8), Inches(4.4), Inches(4.4))
    c2.fill.solid()
    c2.fill.fore_color.rgb = GOLD
    c2.line.color.rgb = DARK_GREEN
    l1 = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(5), Inches(0.5))
    l1.text_frame.text = "TAM: $72B"
    l1.text_frame.paragraphs[0].font.bold = True
    l1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_slide_12_ask(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE_BG)
    add_title(slide, "The Ask: $500,000 Pre-Seed.")
    chart_data = CategoryChartData()
    chart_data.categories = ['Product', 'GTM', 'Team']
    chart_data.add_series('Series 1', (0.4, 0.3, 0.3))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(1), Inches(2), Inches(6), Inches(6), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

def generate_final_reconstruction():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    add_slide_1_cover(prs)
    add_slide_2_stats(prs)
    add_slide_7_market(prs)
    add_slide_12_ask(prs)
    for i in range(13-4):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(s, BEIGE_BG)
        add_title(s, f"Slide {5+i}: Native Reconstruction In Progress")
    
    output_path = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/FINAL_Pitchdeck/AuraGuide_NATIVE_v4_RECONSTRUCTED.pptx"
    prs.save(output_path)
    print(f"Final Native Reconstructed Deck saved to {output_path}")

if __name__ == "__main__":
    generate_final_reconstruction()
