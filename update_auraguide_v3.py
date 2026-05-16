from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import os

# COLORS (Consistent with Viktron Health AuraGuide Branding)
BEIGE_BG = RGBColor(240, 242, 235)
DARK_GREEN = RGBColor(27, 43, 17)
LIGHT_GREEN = RGBColor(225, 234, 216)
GOLD = RGBColor(181, 150, 75)
WHITE = RGBColor(255, 255, 255)

def set_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BEIGE_BG

def add_title(slide, text, size=44):
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(15), Inches(1))
    tf = title.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.name = 'Playfair Display'
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    return title

def add_stat_card(slide, left, top, width, height, number, label):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = DARK_GREEN
    box.line.width = Pt(4)
    n = slide.shapes.add_textbox(left, top + Inches(0.3), width, Inches(1))
    n.text_frame.text = number
    p = n.text_frame.paragraphs[0]
    p.font.name = 'Playfair Display'
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    l = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(1.3), width - Inches(0.2), height - Inches(1.4))
    l.text_frame.text = label
    p = l.text_frame.paragraphs[0]
    p.font.name = 'Outfit'
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER

def build_v3_deck():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # 1. COVER
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    # Circle Decoration
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-4), Inches(-1), Inches(12), Inches(11))
    s.fill.solid()
    s.fill.fore_color.rgb = DARK_GREEN
    s.line.fill.background()
    t = slide.shapes.add_textbox(Inches(8), Inches(3.2), Inches(7), Inches(1.5))
    t.text_frame.text = "AURAGUIDE"
    t.text_frame.paragraphs[0].font.name = 'Playfair Display'
    t.text_frame.paragraphs[0].font.size = Pt(100)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = DARK_GREEN
    sub = slide.shapes.add_textbox(Inches(8), Inches(4.8), Inches(7), Inches(0.8))
    sub.text_frame.text = "A Lifeline for Family Caregivers."
    sub.text_frame.paragraphs[0].font.name = 'Outfit'
    sub.text_frame.paragraphs[0].font.size = Pt(28)
    sub.text_frame.paragraphs[0].font.color.rgb = DARK_GREEN

    # 2. MARKET SCOPE (v3 Content)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "63 Million Invisible Workers. $1 Trillion in Unpaid Labor.")
    add_stat_card(slide, Inches(1), Inches(2), Inches(3.2), Inches(3), "63M", "Unpaid family caregivers in the US, 19% increase since 2020.")
    add_stat_card(slide, Inches(4.5), Inches(2), Inches(3.2), Inches(3), "$1T", "Annual economic value of unpaid caregiving in 2024.")
    add_stat_card(slide, Inches(8), Inches(2), Inches(3.2), Inches(3), "27h/wk", "Avg weekly caregiving hours; 40%+ high-intensity.")
    add_stat_card(slide, Inches(11.5), Inches(2), Inches(3.2), Inches(3), "$10.2K", "Avg annual out-of-pocket spend per caregiver.")

    # 3. THE PROBLEM
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "The Caregiver Collapse Cycle.")
    # Content Columns
    cols = [("Phase 1: Physical Blockade", "90% of medication errors reported occur at home. Hands-full reality."),
            ("Phase 2: Knowledge Gap", "78% manage complex tasks with zero clinical background."),
            ("Phase 3: Critical Errors", "Medication errors cost US system $300B annually."),
            ("Phase 4: Burnout", "87% report stress. ER use rises 73% when caregivers burn out.")]
    for i, (h, b) in enumerate(cols):
        left = Inches(1 + i*3.5)
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2), Inches(3), Inches(4))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GREEN
        box.line.color.rgb = DARK_GREEN
        t = slide.shapes.add_textbox(left, Inches(2.2), Inches(3), Inches(3.5))
        tf = t.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = h
        p.font.bold = True
        p2 = tf.add_paragraph()
        p2.text = b
        p2.font.size = Pt(14)

    # 11. BUSINESS MODEL (v3 Pricing: $14.99)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Software-First, Recurring Revenue Engine.")
    tiers = [("Free", "Basic logging & reminders. Onboarding funnel."),
             ("Premium - $14.99/mo", "Full Voice AI, Clinical Interaction Checking, Wellbeing monitoring."),
             ("Enterprise", "Agencies & Discharge programs. Admin dashboard.")]
    for i, (h, b) in enumerate(tiers):
        add_stat_card(slide, Inches(1 + i*5), Inches(2), Inches(4), Inches(4), h, b)

    # 13. ASK ($500K v3)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "The Ask: $500,000 Pre-Seed.")
    data = CategoryChartData()
    data.categories = ['Product', 'GTM', 'Team']
    data.add_series('Funds', (0.4, 0.3, 0.3))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(1), Inches(2), Inches(6), Inches(6), data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    out = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/FINAL_Pitchdeck/AuraGuide_v3_Updated.pptx"
    prs.save(out)
    print(f"v3 Updated Deck saved to {out}")

if __name__ == "__main__":
    build_v3_deck()
