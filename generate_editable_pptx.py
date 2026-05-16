from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
BEIGE = RGBColor(240, 242, 235)
DARK_GREEN = RGBColor(27, 43, 17)
LIME = RGBColor(204, 255, 0)
WHITE = RGBColor(255, 255, 255)

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_stat_card(slide, left, top, width, height, number, label):
    # Shape
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = DARK_GREEN
    shape.line.width = Pt(1)
    
    # Number
    num_box = slide.shapes.add_textbox(left, top + Inches(0.2), width, Inches(0.6))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.name = 'Playfair Display'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    # Label
    label_box = slide.shapes.add_textbox(left, top + Inches(0.8), width, Inches(0.4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = 'Outfit'
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER

def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(1.5))
    tf = title.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.name = 'Playfair Display'
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(1))
    tf = subtitle.text_frame
    tf.text = subtitle_text
    p = tf.paragraphs[0]
    p.font.name = 'Outfit'
    p.font.size = Pt(28)
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER

def add_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title.text_frame.text = "The Caregiver Collapse Cycle."
    title.text_frame.paragraphs[0].font.name = 'Playfair Display'
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = DARK_GREEN
    
    # 4 Stat Cards
    add_stat_card(slide, Inches(1), Inches(2), Inches(3), Inches(1.5), "63M", "Unpaid Caregivers")
    add_stat_card(slide, Inches(4.5), Inches(2), Inches(3), Inches(1.5), "$1T", "Economic Value")
    add_stat_card(slide, Inches(8), Inches(2), Inches(3), Inches(1.5), "27h/wk", "Avg. Care Hours")
    add_stat_card(slide, Inches(11.5), Inches(2), Inches(3), Inches(1.5), "87%", "Report Stress")

    # Content
    content = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(14), Inches(4))
    tf = content.text_frame
    points = [
        "• Invisible to Tech: Caregivers are the operating system keeping care together.",
        "• The Hands-Full Problem: Active care requires full physical involvement.",
        "• Medical Knowledge Gap: 78% manage complex tasks with zero training.",
        "• Fragmented Tools: Current apps lack clinical context and voice-first design."
    ]
    for pt in points:
        p = tf.add_paragraph()
        p.text = pt
        p.font.name = 'Outfit'
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_GREEN
        p.space_after = Pt(15)

def add_solution_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title.text_frame.text = "AuraGuide: Medical Intelligence Partner."
    title.text_frame.paragraphs[0].font.name = 'Playfair Display'
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = DARK_GREEN
    
    # 4 Features
    features = [
        ("Reduce Cognitive Load", "Voice-first intelligence for schedules and interactions."),
        ("Navigate Urgency", "Crisis intelligence with real-time response protocols."),
        ("Protect Caregiver", "Burnout detection and wellbeing monitoring."),
        ("Bridge Language Gaps", "40+ language support with clinical vocabulary.")
    ]
    
    for i, (head, body) in enumerate(features):
        left = Inches(1 + (i % 2) * 7.5)
        top = Inches(2 + (i // 2) * 2.5)
        
        box = slide.shapes.add_textbox(left, top, Inches(6.5), Inches(2))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = head
        p.font.name = 'Playfair Display'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = DARK_GREEN
        
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.name = 'Outfit'
        p2.font.size = Pt(18)
        p2.font.color.rgb = DARK_GREEN

def generate_editable_deck():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Cover
    add_title_slide(prs, "AURAGUIDE", "A Lifeline for Family Caregivers.")
    
    # Problem
    add_problem_slide(prs)
    
    # Solution
    add_solution_slide(prs)
    
    # Market Timing (Simplified)
    add_content_slide(prs, "A $72B Market Opportunity.", [
        "• TAM: $72B (US Family Caregiver Support Market)",
        "• SAM: $8.5B (Digital Caregiver Tech & Voice AI)",
        "• 10,000 boomers turn 65 every day — demand is locked in.",
        "• Voice AI is mainstream: 98%+ accuracy in noisy environments.",
        "• Complete White Space: No competitor offers proactive clinical intelligence."
    ])
    
    # Business Model
    add_content_slide(prs, "SaaS-First Revenue Engine.", [
        "• Tier 1: Free Freemium Onramp (Reminders, logging)",
        "• Tier 2: Premium - $29/month (Clinical AI, family sharing)",
        "• Tier 3: Enterprise (Facility licensing, admin dashboard)",
        "• LTV: $1,044 per premium subscriber (3-year target)",
        "• Gross Margin: ~85% scalable software infrastructure."
    ])

    # The Ask
    add_content_slide(prs, "The Ask: $500,000 Pre-Seed.", [
        "• 40% Product Hardening: App Store launch & clinical benchmarking.",
        "• 30% GTM & Beta: 15-20 caregiver beta program operations.",
        "• 30% Team & Compliance: Clinical advisor hire & HIPAA groundwork.",
        "• Terms: Post-Money SAFE | $3M–$5M Cap | 20% Discount",
        "• Target Launch: Q2 2026 Production Launch."
    ])
    
    output_path = "/Users/vikashvardhan/IdeaProjects/viktronHealthInc/FINAL_Pitchdeck/AuraGuide_PitchDeck_Editable_v4.pptx"
    prs.save(output_path)
    print(f"Editable deck saved to {output_path}")

def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BEIGE)
    
    # Side Border
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_GREEN
    shape.line.fill.background()
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    tf = title.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.name = 'Playfair Display'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    
    # Content
    content = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(14), Inches(6))
    tf = content.text_frame
    tf.word_wrap = True
    
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.name = 'Outfit'
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_GREEN
        p.space_after = Pt(18)

if __name__ == "__main__":
    generate_editable_deck()
