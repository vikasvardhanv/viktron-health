import fitz
from pptx import Presentation
from pptx.util import Inches
import os

def convert_pdf_to_pptx(pdf_path, output_pptx):
    doc = fitz.open(pdf_path)
    prs = Presentation()
    
    # Use the first page to set slide dimensions
    first_page = doc[0]
    # convert points to inches (1 point = 1/72 inch)
    width = first_page.rect.width / 72
    height = first_page.rect.height / 72
    
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    
    temp_dir = "temp_pdf_frames"
    os.makedirs(temp_dir, exist_ok=True)
    
    for i, page in enumerate(doc):
        # High-res render
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2)) 
        img_path = f"{temp_dir}/page_{i}.png"
        pix.save(img_path)
        
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Add image filling the whole slide
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        print(f"Processed page {i+1}")
        
    prs.save(output_pptx)
    print(f"Saved PPTX to {output_pptx}")

if __name__ == "__main__":
    src = "/Users/vikashvardhan/IdeaProjects/PaperResearch/all data/oldimages/AuraGuide_NotebookLM_v4.pdf"
    dest = "FINAL_Pitchdeck/AuraGuide_PitchDeck_v4.pptx"
    convert_pdf_to_pptx(src, dest)
